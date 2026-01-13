from flask import Flask, request, send_file, render_template
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import requests
from pptx.dml.color import RGBColor

app = Flask(__name__)

# --- Helper Function: Convert Hex (#FFFFFF) to RGB (255, 255, 255) ---
def hex_to_rgb(hex_color):
    """Converts a hex color string to an RGB tuple."""
    try:
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    except ValueError:
        return (255, 255, 255) # Default to white if error

def get_image_stream(url):
    try:
        if "github.com" in url and "/blob/" in url:
            url = url.replace("/blob/", "/raw/")
        
        print(f"Downloading image from: {url}")
        response = requests.get(url, timeout=15)
        response.raise_for_status()
        return BytesIO(response.content)
    except Exception as e:
        print(f"Error downloading image: {e}")
        return None

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload-and-design', methods=['POST'])
def upload_and_design():
    if 'pptFile' not in request.files:
        return "No file part", 400

    file = request.files['pptFile']
    design_url = request.form.get('design_url')
    
    # --- NEW: Get Background Color from Form ---
    bg_color_hex = request.form.get('bg_color') 
    
    if not design_url:
        return "No design selected", 400

    if file:
        try:
            prs = Presentation(file)
            image_stream = get_image_stream(design_url)
            
            if not image_stream:
                return "Failed to download the selected design image.", 500

            # Convert the hex color to RGB
            if bg_color_hex:
                r, g, b = hex_to_rgb(bg_color_hex)

            for slide in prs.slides:
                
                # --- 1. Apply Background Color ---
                if bg_color_hex:
                    try:
                        background = slide.background
                        fill = background.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(r, g, b)
                    except Exception as e:
                        print(f"Could not set background: {e}")

                # --- 2. Add Selected Image ---
                try:
                    image_stream.seek(0)
                    left = Inches(7.0) 
                    top = Inches(0.5)
                    width = Inches(2.5)
                    slide.shapes.add_picture(image_stream, left, top, width=width)
                except Exception as e:
                    print(f"Could not add picture: {e}")

            output = BytesIO()
            prs.save(output)
            output.seek(0)

            return send_file(
                output,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name='designed_presentation.pptx'
            )

        except Exception as e:
            return f"Server Error: {str(e)}", 500

if __name__ == '__main__':
    app.run(debug=True, port=5000)
